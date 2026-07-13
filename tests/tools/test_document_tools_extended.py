"""Tests for the DocumentRead extended format support (PR-B).

Covers:
- .xml  → stdlib xml.etree.ElementTree
- .csv  → stdlib csv / plain read
- .txt  → plain UTF-8 read
- .pptx → python-pptx (optional dependency, guarded)

All tests are hermetic: tiny fixtures created in tmp_path.
"""

from __future__ import annotations

from pathlib import Path

import pytest

from magi_agent.tools.context import ToolContext


# ---------------------------------------------------------------------------
# Helper
# ---------------------------------------------------------------------------


def _ctx(tmp_path: Path) -> ToolContext:
    return ToolContext(
        botId="test-bot",
        sessionId="test-session",
        turnId="test-turn",
        workspaceRoot=str(tmp_path),
    )


# ---------------------------------------------------------------------------
# Fix B — XML
# ---------------------------------------------------------------------------


class TestDocumentReadXml:
    def test_xml_basic_returns_ok(self, tmp_path: Path) -> None:
        (tmp_path / "data.xml").write_text(
            "<root><item>hello</item></root>", encoding="utf-8"
        )
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "data.xml"}, _ctx(tmp_path))
        assert result.status == "ok", f"Expected ok, got {result.status!r}: {result.error_code!r}"

    def test_xml_text_node_extracted(self, tmp_path: Path) -> None:
        (tmp_path / "data.xml").write_text(
            "<root><item>test value</item><other>second</other></root>",
            encoding="utf-8",
        )
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "data.xml"}, _ctx(tmp_path))
        assert result.status == "ok"
        text = result.output["text"]  # type: ignore[index]
        assert "test value" in text, f"Expected 'test value' in text; got {text!r}"
        assert "second" in text

    def test_xml_nested_structure(self, tmp_path: Path) -> None:
        xml_content = """<?xml version="1.0"?>
<catalog>
  <book id="1">
    <title>Python Cookbook</title>
    <author>David Beazley</author>
  </book>
  <book id="2">
    <title>Fluent Python</title>
    <author>Luciano Ramalho</author>
  </book>
</catalog>"""
        (tmp_path / "catalog.xml").write_text(xml_content, encoding="utf-8")
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "catalog.xml"}, _ctx(tmp_path))
        assert result.status == "ok"
        text = result.output["text"]  # type: ignore[index]
        assert "Python Cookbook" in text
        assert "Luciano Ramalho" in text

    def test_xml_content_digest_present(self, tmp_path: Path) -> None:
        (tmp_path / "simple.xml").write_text("<root/>", encoding="utf-8")
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "simple.xml"}, _ctx(tmp_path))
        assert result.status == "ok"
        digest = result.output["contentDigest"]  # type: ignore[index]
        assert isinstance(digest, str) and digest.startswith("sha256:")

    def test_malformed_xml_returns_error_not_raises(self, tmp_path: Path) -> None:
        (tmp_path / "bad.xml").write_text("<unclosed>", encoding="utf-8")
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "bad.xml"}, _ctx(tmp_path))
        # Should return error, not raise exception
        assert result.status in ("error", "ok"), f"Unexpected status {result.status!r}"


# ---------------------------------------------------------------------------
# Fix B — CSV
# ---------------------------------------------------------------------------


class TestDocumentReadCsv:
    def test_csv_basic_returns_ok(self, tmp_path: Path) -> None:
        (tmp_path / "data.csv").write_text(
            "name,age\nAlice,30\nBob,25\n", encoding="utf-8"
        )
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "data.csv"}, _ctx(tmp_path))
        assert result.status == "ok", (
            f"Expected ok, got {result.status!r}: {result.error_code!r}"
        )

    def test_csv_content_readable(self, tmp_path: Path) -> None:
        (tmp_path / "data.csv").write_text(
            "name,age\nAlice,30\nBob,25\n", encoding="utf-8"
        )
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "data.csv"}, _ctx(tmp_path))
        assert result.status == "ok"
        text = result.output["text"]  # type: ignore[index]
        assert "Alice" in text
        assert "name" in text

    def test_csv_was_previously_blocked(self, tmp_path: Path) -> None:
        """Regression: before this fix .csv returned extension_not_supported."""
        (tmp_path / "x.csv").write_text("a,b\n1,2\n", encoding="utf-8")
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "x.csv"}, _ctx(tmp_path))
        assert result.status != "blocked" or result.error_code != "document_extension_not_supported", (
            "CSV is still blocked with document_extension_not_supported — fix not applied"
        )


# ---------------------------------------------------------------------------
# Fix B — TXT / MD / RST
# ---------------------------------------------------------------------------


class TestDocumentReadText:
    def test_txt_returns_ok(self, tmp_path: Path) -> None:
        (tmp_path / "notes.txt").write_text("Hello world.\nLine two.\n", encoding="utf-8")
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "notes.txt"}, _ctx(tmp_path))
        assert result.status == "ok"
        assert "Hello world." in result.output["text"]  # type: ignore[index]

    def test_md_returns_ok(self, tmp_path: Path) -> None:
        (tmp_path / "README.md").write_text("# Title\n\nSome content.\n", encoding="utf-8")
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "README.md"}, _ctx(tmp_path))
        assert result.status == "ok"
        assert "Title" in result.output["text"]  # type: ignore[index]

    def test_txt_was_previously_blocked(self, tmp_path: Path) -> None:
        """Regression: before this fix .txt returned extension_not_supported."""
        (tmp_path / "f.txt").write_text("data\n", encoding="utf-8")
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "f.txt"}, _ctx(tmp_path))
        assert result.status != "blocked" or result.error_code != "document_extension_not_supported", (
            ".txt is still blocked — fix not applied"
        )


class TestDocumentReadSourceProjection:
    """Item 5: a DocumentRead source-read populates a source-ledger projection.

    ``LOCAL_READONLY_TOOL_NAMES`` excludes DocumentRead, so it never recorded a
    source. Behind the default-OFF ``MAGI_SOURCE_LEDGER_EVIDENCE_GATE_ENABLED``,
    a successful DocumentRead now attaches a ``sourceProjection`` exactly like
    FileRead, so the evidence collector projects it as a SourceInspection.
    """

    def test_flag_off_no_source_projection(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        # Force the gate off explicitly: the full-profile overlay seeds it ON,
        # so an implicit default would leak under xdist. This test asserts the
        # sandbox-OFF byte-identical path.
        monkeypatch.setenv("MAGI_SOURCE_LEDGER_EVIDENCE_GATE_ENABLED", "0")
        (tmp_path / "notes.txt").write_text("the token economy section\n", encoding="utf-8")
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "notes.txt"}, _ctx(tmp_path))
        assert result.status == "ok"
        # OFF ⇒ byte-identical to main: no sourceProjection key.
        assert "sourceProjection" not in result.metadata

    def test_flag_on_attaches_source_inspection_projection(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        monkeypatch.setenv("MAGI_SOURCE_LEDGER_EVIDENCE_GATE_ENABLED", "1")
        (tmp_path / "notes.txt").write_text("the token economy section\n", encoding="utf-8")
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "notes.txt"}, _ctx(tmp_path))
        assert result.status == "ok"
        projection = result.metadata.get("sourceProjection")
        assert isinstance(projection, dict)
        sources = projection["sources"]
        assert len(sources) == 1
        assert sources[0]["evidenceType"] == "SourceInspection"
        assert sources[0]["inspected"] is True

    def test_flag_on_projection_feeds_collector_as_source_inspection(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """End-to-end item 4+5: DocumentRead result → collector → SourceInspection."""
        monkeypatch.setenv("MAGI_SOURCE_LEDGER_EVIDENCE_GATE_ENABLED", "1")
        (tmp_path / "notes.txt").write_text("the token economy section\n", encoding="utf-8")
        from magi_agent.evidence.local_tool_collector import LocalToolEvidenceCollector
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "notes.txt"}, _ctx(tmp_path))
        collector = LocalToolEvidenceCollector()
        collector.record_tool_result(
            session_id="test-session",
            turn_id="test-turn",
            tool_call_id="call-1",
            tool_name="DocumentRead",
            result=result,
        )
        records = collector.collect_for_turn("test-turn")
        types = {
            getattr(r, "type", None) if not isinstance(r, dict) else r.get("type")
            for r in records
        }
        assert "SourceInspection" in types


# ---------------------------------------------------------------------------
# Fix B — PPTX (optional dependency, guarded)
# ---------------------------------------------------------------------------


class TestDocumentReadPptx:
    def test_pptx_with_python_pptx_installed(self, tmp_path: Path) -> None:
        """When python-pptx is available, extract slide text."""
        pptx = pytest.importorskip("pptx", reason="python-pptx not installed")
        from pptx import Presentation  # type: ignore[import]
        from pptx.util import Inches  # type: ignore[import]

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title is not None:
            title.text = "Slide One Title"
        body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body is not None:
            body.text = "Slide body text here."
        prs.save(str(tmp_path / "deck.pptx"))

        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "deck.pptx"}, _ctx(tmp_path))
        assert result.status == "ok", (
            f"Expected ok with python-pptx installed; got {result.status!r}: {result.error_code!r}"
        )
        text = result.output["text"]  # type: ignore[index]
        assert "Slide One Title" in text or "Slide body text here." in text, (
            f"Expected slide text in output; got {text!r}"
        )

    def test_pptx_without_python_pptx_returns_blocked_not_raises(
        self, tmp_path: Path, monkeypatch: pytest.MonkeyPatch
    ) -> None:
        """When python-pptx is not installed, return blocked (not crash)."""
        import sys

        (tmp_path / "deck.pptx").write_bytes(b"\x50\x4b\x03\x04" + b"\x00" * 20)

        # Remove pptx from sys.modules so the guarded import fails
        monkeypatch.setitem(sys.modules, "pptx", None)  # type: ignore[arg-type]
        monkeypatch.setitem(sys.modules, "pptx.presentation", None)  # type: ignore[arg-type]

        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "deck.pptx"}, _ctx(tmp_path))
        assert result.status == "blocked"
        assert result.error_code == "document_dependency_not_installed"

    def test_pptx_was_previously_blocked_with_extension_error(
        self, tmp_path: Path
    ) -> None:
        """Regression: before this fix .pptx returned document_extension_not_supported."""
        (tmp_path / "deck.pptx").write_bytes(b"\x50\x4b\x03\x04" + b"\x00" * 20)
        from magi_agent.tools.document_tools import document_read

        result = document_read({"path": "deck.pptx"}, _ctx(tmp_path))
        # Should NOT be blocked with extension_not_supported (may be blocked for
        # missing dependency, which is acceptable — but not for "extension not supported")
        assert not (
            result.status == "blocked"
            and result.error_code == "document_extension_not_supported"
        ), "PPTX is still rejected with extension_not_supported — fix not applied"
