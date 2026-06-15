"""DocumentRead tool — extract text from documents in the workspace.

Requires optional extras (``uv sync --extra files``):
  - PDF:  ``pypdf>=4.0``
  - DOCX: ``python-docx>=1.1``
  - PPTX: ``python-pptx>=1.0``

XML, CSV, TXT, MD, and RST files use stdlib only (no extra package required).

When a required package is not installed the handler returns
``status="blocked"`` with ``errorCode="document_dependency_not_installed"``.
"""

from __future__ import annotations

import hashlib
import json
import re
from collections.abc import Mapping
from pathlib import Path

from .context import ToolContext
from .result import ToolResult
from .spreadsheet_tools import (
    _SpreadsheetPolicyError,
    _base_metadata,
    _blocked_result,
    _error_result,
    _resolve_workspace_path,
    _sanitize_text,
    _workspace_root,
    _markdown_table,
)
from .truncation import cap_text

_MAX_DOCUMENT_BYTES = 20 * 1024 * 1024  # 20 MiB
_DEFAULT_MAX_CHARS = 40_000
_MAX_CHARS = 200_000

_SUPPORTED_EXTENSIONS_CORE = frozenset({".pdf", ".docx"})

# Extended formats: .pptx requires python-pptx (optional); the rest use stdlib only.
_SUPPORTED_EXTENSIONS_EXTENDED = frozenset(
    {".pptx", ".xml", ".csv", ".txt", ".md", ".rst"}
)

_SUPPORTED_EXTENSIONS = _SUPPORTED_EXTENSIONS_CORE | _SUPPORTED_EXTENSIONS_EXTENDED


# ---------------------------------------------------------------------------
# Public handler
# ---------------------------------------------------------------------------


def document_read(arguments: Mapping[str, object], context: ToolContext) -> ToolResult:
    """Extract text from a PDF or DOCX file in the workspace.

    PDF pages are joined with ``\\n\\n---\\n\\n``; DOCX paragraphs and tables
    are rendered as plain text / markdown.  Output is capped at *maxChars*
    (default 40 000) and run through ``_sanitize_text`` before returning.
    """
    tool_name = "document_read"
    path_text = _str_arg(arguments, "path")
    if path_text is None:
        return _blocked_result(tool_name, "path_required")

    try:
        root = _workspace_root(context)
        resolved = _resolve_workspace_path(root, path_text, must_exist=True)
    except _SpreadsheetPolicyError as error:
        return _blocked_result(tool_name, error.reason_code)
    except OSError:
        return _error_result(tool_name, "document_read_failed")

    suffix = Path(resolved.relative).suffix.casefold()
    if suffix not in _SUPPORTED_EXTENSIONS:
        return _blocked_result(
            tool_name,
            "document_extension_not_supported",
            f"Supported extensions: {', '.join(sorted(_SUPPORTED_EXTENSIONS))}",
        )

    try:
        byte_size = resolved.path.stat().st_size
    except OSError:
        return _error_result(tool_name, "document_read_failed")

    if byte_size > _MAX_DOCUMENT_BYTES:
        return _error_result(tool_name, "document_input_too_large")

    max_chars_raw = arguments.get("maxChars")
    max_chars = _DEFAULT_MAX_CHARS
    if isinstance(max_chars_raw, int) and not isinstance(max_chars_raw, bool):
        max_chars = min(max(max_chars_raw, 100), _MAX_CHARS)

    if suffix == ".pdf":
        result = _read_pdf(resolved, arguments, tool_name, max_chars, byte_size)
    elif suffix == ".docx":
        result = _read_docx(resolved, tool_name, max_chars, byte_size)
    elif suffix == ".pptx":
        result = _read_pptx(resolved, tool_name, max_chars, byte_size)
    elif suffix == ".xml":
        result = _read_xml(resolved, tool_name, max_chars, byte_size)
    else:
        # .csv / .txt / .md / .rst — plain UTF-8 text (stdlib only)
        result = _read_text(resolved, tool_name, max_chars, byte_size)
    # A successful DocumentRead is a SOURCE READ (the live model reads documents
    # via DocumentRead, not FileRead). ``LOCAL_READONLY_TOOL_NAMES`` excludes
    # DocumentRead, so it never populated a source ledger and a document-grounded
    # turn could not satisfy the source-evidence gate. Behind the default-OFF
    # ``MAGI_SOURCE_LEDGER_EVIDENCE_GATE_ENABLED`` flag, attach the SAME
    # ``sourceProjection`` a FileRead emits (built via the existing
    # ``LocalResearchSourceLedger.record_source`` + ``public_source_ledger_report``)
    # so the collector projects it as a SourceInspection EvidenceRecord. When the
    # flag is OFF this is a no-op and the result is byte-identical to main.
    return _maybe_attach_source_projection(result, context, resolved)


def _maybe_attach_source_projection(
    result: ToolResult,
    context: ToolContext,
    resolved: object,
) -> ToolResult:
    if result.status != "ok":
        return result
    try:
        import os  # noqa: PLC0415

        from magi_agent.config.env import (  # noqa: PLC0415
            parse_source_ledger_evidence_gate_enabled,
        )

        if not parse_source_ledger_evidence_gate_enabled(os.environ):
            return result
        if "sourceProjection" in result.metadata:
            return result
        from magi_agent.evidence.source_ledger import (  # noqa: PLC0415
            LocalResearchSourceLedger,
            public_source_ledger_report,
        )

        session_ref = _digest_text(context.session_id or "session")
        turn_ref = context.turn_id or "unknown-turn"
        path_ref = getattr(resolved, "path_ref", None) or "document"
        ledger = LocalResearchSourceLedger(
            ledgerId=f"ledger:{_digest_text(str(path_ref))}",
            sessionId=f"session:{session_ref}",
            turnId=turn_ref,
        )
        ledger.record_source(
            {
                "turnId": turn_ref,
                "toolName": "DocumentRead",
                "toolUseId": context.tool_use_id or "DocumentRead:local",
                "evidenceType": "SourceInspection",
                "kind": "external_doc",
                "uri": f"workspace://{path_ref}",
                "inspected": True,
                "contentType": "text/plain",
            }
        )
        projection = public_source_ledger_report(ledger).model_dump(
            by_alias=True, mode="json", warnings=False
        )
        return result.model_copy(
            update={"metadata": {**result.metadata, "sourceProjection": projection}}
        )
    except Exception:
        return result


# ---------------------------------------------------------------------------
# PDF reader
# ---------------------------------------------------------------------------


def _read_pdf(
    resolved: object,
    arguments: Mapping[str, object],
    tool_name: str,
    max_chars: int,
    byte_size: int,
) -> ToolResult:
    from .spreadsheet_tools import _ResolvedPath  # noqa: PLC0415

    assert isinstance(resolved, _ResolvedPath)

    try:
        from pypdf import PdfReader  # noqa: PLC0415
    except ImportError:
        return _blocked_result(tool_name, "document_dependency_not_installed")

    page_range_str = _str_arg(arguments, "pageRange")

    try:
        reader = PdfReader(str(resolved.path))
        total_pages = len(reader.pages)
        start_page, end_page = _parse_page_range(page_range_str, total_pages)

        parts: list[str] = []
        for page_idx in range(start_page, end_page):
            text = reader.pages[page_idx].extract_text() or ""
            parts.append(text)
        raw_text = "\n\n---\n\n".join(parts)
    except Exception:  # noqa: BLE001
        return _error_result(tool_name, "pdf_read_error")

    return _finish_document_result(
        tool_name,
        raw_text,
        page_count=total_pages,
        byte_size=byte_size,
        max_chars=max_chars,
        path_ref=resolved.path_ref,
        content_digest=_digest_path(resolved.path),
    )


# ---------------------------------------------------------------------------
# DOCX reader
# ---------------------------------------------------------------------------


def _read_docx(
    resolved: object,
    tool_name: str,
    max_chars: int,
    byte_size: int,
) -> ToolResult:
    from .spreadsheet_tools import _ResolvedPath  # noqa: PLC0415

    assert isinstance(resolved, _ResolvedPath)

    try:
        from docx import Document  # noqa: PLC0415
    except ImportError:
        return _blocked_result(tool_name, "document_dependency_not_installed")

    try:
        doc = Document(str(resolved.path))
        raw_text = extract_docx_text(doc)
    except Exception:  # noqa: BLE001
        return _error_result(tool_name, "docx_read_error")

    return _finish_document_result(
        tool_name,
        raw_text,
        page_count=None,
        byte_size=byte_size,
        max_chars=max_chars,
        path_ref=resolved.path_ref,
        content_digest=_digest_path(resolved.path),
    )


def extract_docx_text(doc: object) -> str:
    """Extract paragraph + table text from a ``python-docx`` ``Document``.

    Walks the body elements in order (paragraphs + tables interleaved) and
    renders tables as markdown — the same logic ``document_read`` uses, factored
    out so other surfaces (e.g. the DocumentWrite coverage verifier) can extract
    the rendered text from an in-memory ``Document`` without going through disk.

    The ``docx`` import stays lazy (inside this function) so importing this
    module never pulls ``docx`` into ``sys.modules``; callers already hold a
    constructed ``Document`` instance.
    """
    from docx.oxml.ns import qn  # noqa: PLC0415

    parts: list[str] = []
    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag
        if tag == "p":
            text_nodes = element.findall(f".//{qn('w:t')}")
            para_text = "".join(t.text or "" for t in text_nodes)
            if para_text.strip():
                parts.append(para_text)
        elif tag == "tbl":
            rows: list[list[str]] = []
            for tr in element.findall(f".//{qn('w:tr')}"):
                row_cells: list[str] = []
                for tc in tr.findall(f".//{qn('w:tc')}"):
                    cell_texts = tc.findall(f".//{qn('w:t')}")
                    row_cells.append("".join(t.text or "" for t in cell_texts))
                if row_cells:
                    rows.append(row_cells)
            if rows:
                width = max(len(r) for r in rows)
                padded = [r + [""] * (width - len(r)) for r in rows]
                parts.append(_markdown_table(padded))

    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# PPTX reader
# ---------------------------------------------------------------------------


def _read_pptx(
    resolved: object,
    tool_name: str,
    max_chars: int,
    byte_size: int,
) -> ToolResult:
    from .spreadsheet_tools import _ResolvedPath  # noqa: PLC0415

    assert isinstance(resolved, _ResolvedPath)

    try:
        from pptx import Presentation  # noqa: PLC0415
    except ImportError:
        return _blocked_result(tool_name, "document_dependency_not_installed")

    try:
        prs = Presentation(str(resolved.path))
        parts: list[str] = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_texts.append(line)
            if slide_texts:
                parts.append(f"[Slide {slide_idx}]\n" + "\n".join(slide_texts))
        raw_text = "\n\n".join(parts)
    except Exception:  # noqa: BLE001
        return _error_result(tool_name, "pptx_read_error")

    return _finish_document_result(
        tool_name,
        raw_text,
        page_count=len(prs.slides),
        byte_size=byte_size,
        max_chars=max_chars,
        path_ref=resolved.path_ref,
        content_digest=_digest_path(resolved.path),
    )


# ---------------------------------------------------------------------------
# XML reader
# ---------------------------------------------------------------------------


def _read_xml(
    resolved: object,
    tool_name: str,
    max_chars: int,
    byte_size: int,
) -> ToolResult:
    """Extract text nodes from an XML file using stdlib xml.etree.ElementTree."""
    import xml.etree.ElementTree as ET  # noqa: PLC0415

    from .spreadsheet_tools import _ResolvedPath  # noqa: PLC0415

    assert isinstance(resolved, _ResolvedPath)

    try:
        tree = ET.parse(str(resolved.path))  # noqa: S314
        root = tree.getroot()
        texts: list[str] = []

        def _collect(element: ET.Element) -> None:
            # Text content of the element itself
            if element.text and element.text.strip():
                texts.append(element.text.strip())
            # Tail text (text after the closing tag, before the next sibling tag)
            if element.tail and element.tail.strip():
                texts.append(element.tail.strip())
            for child in element:
                _collect(child)

        _collect(root)
        raw_text = "\n".join(texts)
    except ET.ParseError:
        return _error_result(tool_name, "xml_parse_error")
    except Exception:  # noqa: BLE001
        return _error_result(tool_name, "xml_read_error")

    return _finish_document_result(
        tool_name,
        raw_text,
        page_count=None,
        byte_size=byte_size,
        max_chars=max_chars,
        path_ref=resolved.path_ref,
        content_digest=_digest_path(resolved.path),
    )


# ---------------------------------------------------------------------------
# Plain-text reader (.csv / .txt / .md / .rst)
# ---------------------------------------------------------------------------


def _read_text(
    resolved: object,
    tool_name: str,
    max_chars: int,
    byte_size: int,
) -> ToolResult:
    """Read a plain-text file (CSV, TXT, MD, RST) as UTF-8."""
    from .spreadsheet_tools import _ResolvedPath  # noqa: PLC0415

    assert isinstance(resolved, _ResolvedPath)

    try:
        raw_text = resolved.path.read_text(encoding="utf-8", errors="replace")
    except Exception:  # noqa: BLE001
        return _error_result(tool_name, "text_read_error")

    return _finish_document_result(
        tool_name,
        raw_text,
        page_count=None,
        byte_size=byte_size,
        max_chars=max_chars,
        path_ref=resolved.path_ref,
        content_digest=_digest_path(resolved.path),
    )


# ---------------------------------------------------------------------------
# Shared finish
# ---------------------------------------------------------------------------


def _finish_document_result(
    tool_name: str,
    raw_text: str,
    *,
    page_count: int | None,
    byte_size: int,
    max_chars: int,
    path_ref: str,
    content_digest: str,
) -> ToolResult:
    sanitized, redacted = _sanitize_text(raw_text)
    sanitized, truncated = cap_text(sanitized, max_chars)

    output: dict[str, object] = {
        "text": sanitized,
        "truncated": truncated,
        "contentDigest": content_digest,
    }
    if page_count is not None:
        output["pageCount"] = page_count
    if redacted:
        output["redacted"] = True

    return ToolResult(
        status="ok",
        output=output,
        llmOutput=output,
        transcriptOutput={
            "toolName": tool_name,
            "charCount": len(sanitized),
            "truncated": truncated,
            "contentDigest": content_digest,
        },
        metadata={
            **_base_metadata(tool_name, permission_class="read", mutates_workspace=False),
            "contentDigest": content_digest,
            "byteCount": byte_size,
            "charCount": len(sanitized),
            "truncated": truncated,
            "pathRef": path_ref,
            "redactionStatus": "redacted" if redacted else "no_redaction_needed",
        },
    )


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _parse_page_range(page_range_str: str | None, total_pages: int) -> tuple[int, int]:
    """Parse a ``pageRange`` string like ``'1-5'`` or ``'3'`` to 0-based indices.

    Returns ``(start, end)`` suitable for ``range(start, end)``.
    """
    if not page_range_str:
        return 0, total_pages
    page_range_str = page_range_str.strip()
    m_range = re.fullmatch(r"(\d+)-(\d+)", page_range_str)
    if m_range:
        start = max(int(m_range.group(1)) - 1, 0)
        end = min(int(m_range.group(2)), total_pages)
        return start, end
    m_single = re.fullmatch(r"(\d+)", page_range_str)
    if m_single:
        page = max(int(m_single.group(1)) - 1, 0)
        return page, min(page + 1, total_pages)
    return 0, total_pages


def _str_arg(arguments: Mapping[str, object], name: str) -> str | None:
    value = arguments.get(name)
    if isinstance(value, str):
        return value
    return None


def _digest_path(path: Path) -> str:
    digest = hashlib.sha256(path.read_bytes()).hexdigest()
    return f"sha256:{digest}"


def _digest_text(value: str) -> str:
    return hashlib.sha256(value.encode("utf-8")).hexdigest()[:24]


# ---------------------------------------------------------------------------
# document_search — in-document term search (PDF only, page-addressed)
# ---------------------------------------------------------------------------

_DOCUMENT_SEARCH_SUPPORTED_EXTENSIONS = frozenset({".pdf"})
_SEARCH_SNIPPET_CONTEXT = 120  # chars on each side of the match
_MAX_SEARCH_MATCHES = 200


def document_search(
    arguments: Mapping[str, object], context: ToolContext
) -> ToolResult:
    """Search a PDF document for a query term, returning page number + snippet
    for each match.

    Only PDF files are supported (text-layer search via pypdf).  Returns
    ``status="blocked"`` with ``errorCode="document_search_not_supported_for_format"``
    for other document types.

    Output fields:
    - ``matches``: list of ``{page: int, snippet: str}`` (1-based page numbers)
    - ``matchCount``: total number of matches found
    - ``totalPages``: total pages in the document
    """
    tool_name = "document_search"

    path_text = _str_arg(arguments, "path")
    if path_text is None:
        return _blocked_result(tool_name, "path_required")

    query_text = _str_arg(arguments, "query")
    if query_text is None:
        return _blocked_result(tool_name, "query_required")

    try:
        root = _workspace_root(context)
        resolved = _resolve_workspace_path(root, path_text, must_exist=True)
    except _SpreadsheetPolicyError as error:
        return _blocked_result(tool_name, error.reason_code)
    except OSError:
        return _error_result(tool_name, "document_search_failed")

    suffix = Path(resolved.relative).suffix.casefold()
    if suffix not in _DOCUMENT_SEARCH_SUPPORTED_EXTENSIONS:
        return _blocked_result(
            tool_name,
            "document_search_not_supported_for_format",
            f"document_search supports: {', '.join(sorted(_DOCUMENT_SEARCH_SUPPORTED_EXTENSIONS))}",
        )

    try:
        byte_size = resolved.path.stat().st_size
    except OSError:
        return _error_result(tool_name, "document_search_failed")

    if byte_size > _MAX_DOCUMENT_BYTES:
        return _error_result(tool_name, "document_input_too_large")

    try:
        from pypdf import PdfReader  # noqa: PLC0415
    except ImportError:
        return _blocked_result(tool_name, "document_dependency_not_installed")

    try:
        reader = PdfReader(str(resolved.path))
        total_pages = len(reader.pages)
        query_lower = query_text.casefold()
        matches: list[dict[str, object]] = []

        for page_idx in range(total_pages):
            page_text = reader.pages[page_idx].extract_text() or ""
            page_lower = page_text.casefold()
            pos = 0
            while True:
                found = page_lower.find(query_lower, pos)
                if found == -1:
                    break
                start = max(0, found - _SEARCH_SNIPPET_CONTEXT)
                end = min(len(page_text), found + len(query_text) + _SEARCH_SNIPPET_CONTEXT)
                snippet = page_text[start:end].replace("\n", " ").strip()
                matches.append({
                    "page": page_idx + 1,  # 1-based
                    "snippet": snippet,
                })
                pos = found + len(query_lower)
                if len(matches) >= _MAX_SEARCH_MATCHES:
                    break
            if len(matches) >= _MAX_SEARCH_MATCHES:
                break

    except Exception:  # noqa: BLE001
        return _error_result(tool_name, "pdf_read_error")

    output: dict[str, object] = {
        "matches": matches,
        "matchCount": len(matches),
        "totalPages": total_pages,
        "query": query_text,
    }
    content_digest = _digest_path(resolved.path)
    return ToolResult(
        status="ok",
        output=output,
        llmOutput=output,
        transcriptOutput={
            "toolName": tool_name,
            "matchCount": len(matches),
            "totalPages": total_pages,
            "contentDigest": content_digest,
        },
        metadata={
            **_base_metadata(tool_name, permission_class="read", mutates_workspace=False),
            "contentDigest": content_digest,
            "byteCount": byte_size,
            "matchCount": len(matches),
            "totalPages": total_pages,
            "pathRef": resolved.path_ref,
        },
    )


__all__ = ["document_read", "document_search", "extract_docx_text"]
